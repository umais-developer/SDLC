"""
Build the SDLC pipeline walkthrough deck.

Generates `docs/SDLC-pipeline-walkthrough.pptx` from the slide content defined
inline below. Re-run this after editing the content; do not edit the .pptx
directly so changes stay reviewable in diffs.

    python docs/build_walkthrough_pptx.py
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

OUT = Path(__file__).resolve().parent / "SDLC-pipeline-walkthrough.pptx"

# --- theme tokens ------------------------------------------------------------
BG = RGBColor(0x0B, 0x0E, 0x14)
PANEL = RGBColor(0x11, 0x16, 0x1F)
FG = RGBColor(0xE6, 0xED, 0xF3)
MUTED = RGBColor(0x9A, 0xA6, 0xB2)
ACCENT = RGBColor(0x43, 0xE8, 0xB8)
PINK = RGBColor(0xFF, 0x5D, 0x8F)
FOCUS = RGBColor(0xFF, 0xD1, 0x66)
BORDER = RGBColor(0x2A, 0x33, 0x42)

WIDTH = Inches(13.333)
HEIGHT = Inches(7.5)

# --- helpers -----------------------------------------------------------------

def set_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    text: str,
    *,
    left: Inches,
    top: Inches,
    width: Inches,
    height: Inches,
    size: int = 18,
    color: RGBColor = FG,
    bold: bool = False,
    align: int = PP_ALIGN.LEFT,
) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Segoe UI"


def add_bullets(
    slide,
    bullets: list[str],
    *,
    left: Inches,
    top: Inches,
    width: Inches,
    height: Inches,
    size: int = 18,
    color: RGBColor = FG,
) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {line}" if not line.startswith(" ") else line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Segoe UI"


def add_accent_bar(slide, top: Inches = Inches(0.6)) -> None:
    """A thin gradient-feel bar at the top of each content slide."""
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), top, Inches(1.2), Inches(0.06))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT


def add_footer(slide, text: str) -> None:
    add_text(
        slide,
        text,
        left=Inches(0.6),
        top=Inches(7.05),
        width=Inches(12),
        height=Inches(0.3),
        size=10,
        color=MUTED,
    )


def title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)

    # Big gradient-ish title (PowerPoint can't do CSS gradients on text without XML
    # surgery — accent the title color instead)
    add_text(
        s, title,
        left=Inches(0.8), top=Inches(2.4), width=Inches(11.7), height=Inches(1.6),
        size=54, color=ACCENT, bold=True,
    )
    add_text(
        s, subtitle,
        left=Inches(0.8), top=Inches(4.0), width=Inches(11.7), height=Inches(0.8),
        size=22, color=FG,
    )
    add_text(
        s, "Internal walkthrough — SDLC pipeline",
        left=Inches(0.8), top=Inches(6.6), width=Inches(11.7), height=Inches(0.6),
        size=14, color=MUTED,
    )


def content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str] | None = None,
    *,
    body_text: str | None = None,
    footer: str = "",
    bullet_size: int = 18,
) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_accent_bar(s)
    add_text(
        s, title,
        left=Inches(0.6), top=Inches(0.78),
        width=Inches(12), height=Inches(0.7),
        size=30, color=FG, bold=True,
    )
    if bullets:
        add_bullets(
            s, bullets,
            left=Inches(0.7), top=Inches(1.7),
            width=Inches(12), height=Inches(5),
            size=bullet_size,
        )
    if body_text:
        add_text(
            s, body_text,
            left=Inches(0.7), top=Inches(1.7),
            width=Inches(12), height=Inches(5),
            size=bullet_size,
        )
    if footer:
        add_footer(s, footer)


def two_column_slide(
    prs: Presentation,
    title: str,
    left_title: str, left_bullets: list[str],
    right_title: str, right_bullets: list[str],
    *, footer: str = "",
) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_accent_bar(s)
    add_text(
        s, title,
        left=Inches(0.6), top=Inches(0.78),
        width=Inches(12), height=Inches(0.7),
        size=30, color=FG, bold=True,
    )
    add_text(s, left_title, left=Inches(0.7), top=Inches(1.7), width=Inches(6), height=Inches(0.5),
             size=20, color=ACCENT, bold=True)
    add_bullets(s, left_bullets, left=Inches(0.7), top=Inches(2.2), width=Inches(6), height=Inches(5),
                size=16)
    add_text(s, right_title, left=Inches(7.0), top=Inches(1.7), width=Inches(6), height=Inches(0.5),
             size=20, color=PINK, bold=True)
    add_bullets(s, right_bullets, left=Inches(7.0), top=Inches(2.2), width=Inches(6), height=Inches(5),
                size=16)
    if footer:
        add_footer(s, footer)


def section_divider(prs: Presentation, label: str, subtitle: str = "") -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_text(
        s, label,
        left=Inches(0.8), top=Inches(2.8), width=Inches(11.7), height=Inches(1.4),
        size=44, color=PINK, bold=True,
    )
    if subtitle:
        add_text(
            s, subtitle,
            left=Inches(0.8), top=Inches(4.2), width=Inches(11.7), height=Inches(0.8),
            size=18, color=MUTED,
        )


def table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    *, footer: str = "",
    col_widths: list[Inches] | None = None,
) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    add_accent_bar(s)
    add_text(
        s, title,
        left=Inches(0.6), top=Inches(0.78),
        width=Inches(12), height=Inches(0.7),
        size=30, color=FG, bold=True,
    )
    n_cols = len(headers)
    n_rows = len(rows) + 1
    table = s.shapes.add_table(
        n_rows, n_cols,
        Inches(0.6), Inches(1.7),
        Inches(12.1), Inches(min(5.0, 0.45 * n_rows + 0.5)),
    ).table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        run.font.bold = True
        run.font.size = Pt(14)
        run.font.color.rgb = ACCENT
        run.font.name = "Segoe UI"

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.size = Pt(13)
            run.font.color.rgb = FG
            run.font.name = "Segoe UI"
    if footer:
        add_footer(s, footer)


# --- deck content ------------------------------------------------------------

def build() -> None:
    prs = Presentation()
    prs.slide_width = WIDTH
    prs.slide_height = HEIGHT

    # 1. Title
    title_slide(
        prs,
        "An SDLC Pipeline for AI Agents",
        "From a one-line prompt to a deployable build, with traceability the whole way.",
    )

    # 2. Agenda
    content_slide(
        prs,
        "Agenda",
        bullets=[
            "Why a structured pipeline (and why we built two of them)",
            "The 9-stage pipeline at a glance",
            "Stage-by-stage walkthrough",
            "Case study: building a Snake game with replays end-to-end",
            "The 3-stage 'lite' pipeline — when full ceremony is overkill",
            "Side-by-side comparison",
            "Retrospective: what went well, what didn't, what we learned",
            "How to teach this to the next person",
        ],
        footer="Estimated walkthrough time: 25-30 minutes",
    )

    # 3. Why
    content_slide(
        prs,
        "Why a structured pipeline?",
        bullets=[
            "AI agents are fast at writing code, slow at deciding what to write.",
            "Without structure, prompts drift: same input -> different output, no traceability.",
            "We wanted: explicit upstream IDs, gated handoffs, and verifiable evidence at every stage.",
            "Goal: make 'agent did the work' indistinguishable in audit from 'a deliberate team did the work'.",
            "Stretch: same skill files run under Claude Code, GitHub Copilot, or any agent that can read Markdown.",
        ],
    )

    # 4. The 9-stage flow
    two_column_slide(
        prs,
        "The 9-stage pipeline",
        "Stages 1-5: define the work",
        [
            "1. PRD — problem.json, goals.json (FR/NFR/CON IDs)",
            "2. Architecture — components, tech stack, layer rules",
            "3. UX — flows, states, error paths, accessibility",
            "4. Epics & Stories — story IDs + traceability matrix",
            "5. Plan — task IDs, file paths, tests per task",
        ],
        "Stages 6-9: build and verify",
        [
            "6. Implementation — code under src/, build + test logs",
            "7. Code Review — verdict APPROVE / CHANGES_REQUIRED",
            "8. UAT — test plan + executed results + bugs",
            "9. Deploy — GitHub Pages config (opt-in, manual)",
            "Each stage halts on its verifier failure.",
        ],
        footer="Every output cites upstream IDs. Verifiers enforce that, structurally.",
    )

    # 5. Skill anatomy
    content_slide(
        prs,
        "Anatomy of a stage skill",
        bullets=[
            "SKILL.md — instructions for the agent (deterministic + LLM steps).",
            "prompts/*.md — the actual prompts the agent runs, with prompt_version.",
            "verify/*.py — deterministic gate. Exit 0 = pass.",
            "Final artifact = JSON (machine) + Markdown (human).",
            "meta block on every JSON: prompt versions + source-hash chain.",
            "Result: drift between an artifact and its inputs is detectable mechanically.",
        ],
        footer="Skill files live under .agents/skills/<stage-name>/.",
    )

    # 6. Case study setup
    section_divider(
        prs,
        "Case Study: Snake with Replays",
        "Greenfield, large size, end-to-end run on the 9-stage pipeline.",
    )

    # 7. Case study spec
    content_slide(
        prs,
        "What the user asked for",
        body_text=(
            '"I want to create a really complete snake game with slick UI and snake getting bigger '
            "and also replay mode with ability to save replays and watch them — play them, change "
            "speeds and rewind them and just give an awesome user experience.\""
        ),
        bullet_size=20,
        footer="One sentence in. The pipeline turns it into a working build.",
    )

    # 8. Case study artifacts
    table_slide(
        prs,
        "What came out (Stages 1-8)",
        ["Stage", "Output", "Key signal"],
        [
            ["1 PRD", "5 goals, 10 FRs, 5 NFRs, 4 CONs", "12 explicit assumptions, all tagged"],
            ["2 Architecture", "21 components, 6 layers", "Engine has zero DOM imports"],
            ["3 UX", "8 flows, 17 UI states", "Accessibility tied to NFR IDs"],
            ["4 Epics & Stories", "6 epics, 25 stories, traceability matrix", "Every P0 FR has >=1 story"],
            ["5 Plan", "42 tasks across 15 dependency tiers", "4 vertical slices identified"],
            ["6 Implement", "44 source files, 22 test files", "112/112 tests pass; 40 KB gzip bundle"],
            ["7 Review", "Verdict: APPROVE", "0 Critical/High; 4 Low follow-ups"],
            ["8 UAT", "101 P0 cases, gate APPROVED", "App health: HTTP 200 on vite preview"],
        ],
        footer="All gates green. End-to-end, hands-off-after-the-prompt.",
    )

    # 9. The win that came from the pipeline
    content_slide(
        prs,
        "The architecture chose itself",
        bullets=[
            "CON-3 (engine must be deterministic) was made explicit in Stage 1.",
            "Stage 2 picked a seeded mulberry32 PRNG + serializable engine state — directly answering CON-3.",
            "Stage 5 mapped the seek requirement (FR-7, <=300 ms scrub) to a snapshot+fast-forward strategy.",
            "Stage 6 implemented it with a 1000-tick determinism fixture as a regression gate.",
            "None of this is exotic engineering. The pipeline made it inevitable instead of optional.",
        ],
        footer="When the constraints are explicit, the design falls out.",
    )

    # 10. Lite pipeline
    section_divider(
        prs,
        "When ceremony is overkill",
        "Introducing /lite-build — the 3-stage variant.",
    )

    # 11. Lite spec
    content_slide(
        prs,
        "/lite-build at a glance",
        bullets=[
            "Three stages: Spec -> Build -> Verify.",
            "One markdown per stage. No JSON, no separate UX flows, no traceability matrix.",
            "Single combined verifier (~150 lines of Python).",
            "Kept: anti-hallucination on assumptions, FR-id traceability into tests, build + tests must pass.",
            "Dropped: stories, tasks, drift hashes, dimension-tagged review, separate UAT.",
            "Use it when: small feature, prototype, focused enhancement.",
        ],
        footer=".agents/skills-lite/  +  /lite-build slash command",
    )

    # 12. Side-by-side
    table_slide(
        prs,
        "Lite vs full — Snake rebuilt under both",
        ["Metric", "9-stage", "3-stage lite"],
        [
            ["Stages", "8 (9 with deploy)", "3"],
            ["Skill files", "9 + ~12 prompts + verifiers", "4 + 1 verifier"],
            ["Pipeline artifacts", "17 JSON + Markdown", "3 files"],
            ["Source files", "21 TS + 1 CSS", "6 TS + 1 CSS"],
            ["Lines of source", "~5000", "~1500"],
            ["Tests", "112 (22 files)", "21 (4 files)"],
            ["Bundle (gzip)", "12.65 KB JS + 2 KB CSS", "8.41 KB JS + 1.43 KB CSS"],
            ["Wall time", "~hours", "~25 minutes"],
            ["Verifier fights (mechanical)", "~12", "1"],
        ],
        footer="Same product, same gates green. Lite trades audit depth for speed.",
    )

    # 13. Section: retro
    section_divider(prs, "Retrospective", "Honest reflection from running both pipelines.")

    # 14. What went well
    content_slide(
        prs,
        "What went well",
        bullets=[
            "Forcing 'large' classification up front cascaded sensibly through every stage.",
            "Traceability chain caught real omissions (missing flow links, untested FRs).",
            "Anti-hallucination rule pushed assumptions out of comments and into a tagged list.",
            "Determinism fixture caught two would-be regressions during Stage 6 retries.",
            "Final deliverable shipped: 40 KB gzip bundle, 200 OK on preview, 112/112 tests.",
            "The lite pipeline got to a working build in roughly 1/3 the time and 1/3 the LOC.",
        ],
    )

    # 15. What we learned
    content_slide(
        prs,
        "What we learned",
        bullets=[
            "Verifiers should reject vagueness, not just structure. The 'NFR threshold' rejection was the most useful LLM-side signal we got.",
            "When two artifacts share signal (problem.json + goals.json), pick one.",
            "Single-source markdown per stage is easier to keep accurate than per-stage JSON+markdown duplication.",
            "Encoding bugs in helpers (em-dashes) silently mangle artifacts on Windows. UTF-8 in/out, every script.",
            "The model's first instinct is to be exhaustive. The pipeline's job is to make it be specific.",
            "Retrospect honestly: 'all tests passed' is necessary, not sufficient. Self-written tests can self-confirm.",
        ],
    )

    # 16. What did not go well
    content_slide(
        prs,
        "What didn't go so well",
        bullets=[
            "Stage 8 was theatre. 101 'test cases' all PASS by pointing at the same Vitest log; structure was checked, substance wasn't.",
            "Stage 7 was self-review — same agent wrote the code and graded it. 4 Low findings logged; a stranger reviewer would've caught more.",
            "Story-coverage comments added to test files to satisfy the verifier are exactly the fake compliance the rule is meant to prevent.",
            "12+ verifier cycles were spent on schema mechanics (line-range bounds, missing 'quality' dimension, key renames). Zero product-quality signal.",
            "inject_meta.py mangled em-dashes into garbled escape sequences — forced a full ASCII rewrite of artifacts.",
            "Two near-identical engines exist (full + lite) because greenfield runs in the same session aren't independent.",
        ],
    )

    # 17. What we could do better
    content_slide(
        prs,
        "What we could do better",
        bullets=[
            "Stage 8: actually run Playwright. Install once, headless Chromium, real screenshots/video/trace as evidence.",
            "Stage 7: run review in a fresh agent with no memory of writing the code. True independence.",
            "Verifiers: stop asking 'does the file mention the ID?' and start asking 'does the test exercise the behavior?'.",
            "Helpers: standardize on UTF-8 + `--ascii-only` mode where Windows terminals can't decode.",
            "Retire duplicate prose: prefer one 'final.md' per stage and treat JSON as machine cache only.",
            "Add a 'budget' to Stage 8 — if PASS rate >=N% with browser evidence, pass; if not, halt with the failed cases.",
        ],
    )

    # 18. How to teach this
    content_slide(
        prs,
        "Teaching it to the next person",
        bullets=[
            "Day 1: read README + STAGE-CONVENTIONS.md. Run /lite-build on a tiny feature. Read every artifact.",
            "Day 2: run /create-product on a small medium-size feature. Stop at each gate. Inspect the verifier output.",
            "Day 3: deliberately break something — change goals.json after stories.json is generated. Run check_drift.py. Watch it catch you.",
            "Week 1: pair with someone running a real customer feature end-to-end. Discuss every assumption tag.",
            "Mental model: 'the pipeline is a contract negotiator.' Stages are not phases of work, they are checkpoints where the agent has to commit to a story upstream.",
            "Failure mode to flag early: passing tests do not equal passing UAT. The two should rhyme but never collapse into each other.",
        ],
    )

    # 19. Quick reference
    table_slide(
        prs,
        "Quick reference",
        ["Want to...", "Do this"],
        [
            ["Build a small feature", "/lite-build <description>"],
            ["Build a product end-to-end", "/create-product <description>"],
            ["Re-run a single stage", "/stage-N (full) or python .agents/skills-lite/_verify.py <stage> (lite)"],
            ["Smoke-test all verifiers", "python .agents/tests/run_verifiers.py"],
            ["Detect upstream/downstream drift", "python .agents/tests/check_drift.py"],
            ["Recover from wedged Stage 6 run", "python .agents/skills/stage-6-implement/recovery.py --status"],
            ["Read the conventions", ".agents/skills/STAGE-CONVENTIONS.md"],
        ],
        col_widths=[Inches(4.5), Inches(8.0)],
    )

    # 20. Closing
    title_slide(
        prs,
        "Thanks.",
        "Questions, war stories, suggestions for the next iteration?",
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT.relative_to(Path.cwd()) if OUT.is_relative_to(Path.cwd()) else OUT}")


if __name__ == "__main__":
    build()
